#!/usr/bin/env python3
"""
Screenshot-based PPTX exporter for the Jina AI + EIS presentation.

Renders each <section> in index.html with a headless browser,
screenshots it at 1920x1080, and places each screenshot as a
full-bleed slide in a PPTX file.

Usage:
    cd docs/presentation
    pip install python-pptx playwright Pillow
    python -m playwright install chromium
    python export_pptx_screenshot.py
    # Output: output/jina-eis-101-screenshots.pptx

Options:
    --html FILE     HTML file to export (default: index.html)
    --output FILE   Output PPTX path (default: output/jina-eis-101-screenshots.pptx)
    --width N       Viewport width (default: 1920)
    --height N      Viewport height (default: 1080)
    --live-only     Only export live deck slides (skip reference appendix)
"""

import argparse
import sys
import tempfile
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu


SCRIPT_DIR = Path(__file__).parent
DEFAULT_HTML = SCRIPT_DIR / "index.html"
DEFAULT_OUTPUT = SCRIPT_DIR / "output" / "jina-eis-101-screenshots.pptx"


def export(html_path: Path, output_path: Path, width: int, height: int,
           live_only: bool):
    html_path = html_path.resolve()
    if not html_path.exists():
        print(f"Error: {html_path} not found")
        sys.exit(1)

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Slide dimensions: 13.333 x 7.5 inches (widescreen 16:9)
    slide_width = Inches(13.333)
    slide_height = Inches(7.5)

    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    blank_layout = prs.slide_layouts[6]  # blank

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": width, "height": height},
                                device_scale_factor=2)
        page.goto(f"file://{html_path}", wait_until="networkidle")

        # Wait for fonts and fade-in animations
        page.wait_for_timeout(1000)

        # Force all fade-up elements to be visible
        page.evaluate("""
            document.querySelectorAll('.fade-up').forEach(el => {
                el.classList.add('visible');
                el.style.opacity = '1';
                el.style.transform = 'translateY(0)';
            });
        """)

        # Get all sections
        sections = page.evaluate("""
            () => {
                const sections = document.querySelectorAll('section[data-nav]');
                return Array.from(sections).map((s, i) => ({
                    index: i,
                    nav: s.dataset.nav,
                    group: s.dataset.group || '',
                }));
            }
        """)

        if live_only:
            sections = [s for s in sections if s["group"] != "reference"]

        total = len(sections)
        print(f"Found {total} sections to export")

        with tempfile.TemporaryDirectory() as tmpdir:
            for i, sec in enumerate(sections):
                label = f"[{i+1}/{total}] {sec['nav']}"
                if sec["group"]:
                    label += f" ({sec['group']})"
                print(f"  Capturing {label}...")

                # Scroll section into view and screenshot it
                img_path = Path(tmpdir) / f"slide_{i:03d}.png"
                page.evaluate(f"""
                    () => {{
                        const sections = document.querySelectorAll('section[data-nav]');
                        const sec = sections[{sec['index']}];
                        sec.scrollIntoView({{ behavior: 'instant' }});
                    }}
                """)
                page.wait_for_timeout(300)

                # Click interactive demo buttons if present in this section
                demo_btn = page.evaluate(f"""
                    () => {{
                        const sections = document.querySelectorAll('section[data-nav]');
                        const sec = sections[{sec['index']}];
                        const btn = sec.querySelector('#rerankBtn, #readerBtn, #vlmBtn');
                        if (btn && btn.offsetParent !== null) {{
                            btn.click();
                            return btn.id;
                        }}
                        return null;
                    }}
                """)
                if demo_btn:
                    # Wait for animation to complete
                    wait_ms = 4000 if demo_btn == "vlmBtn" else 2500
                    print(f"    → Clicked {demo_btn}, waiting {wait_ms}ms for output...")
                    page.wait_for_timeout(wait_ms)

                # Screenshot the viewport (what you'd see on screen)
                page.screenshot(path=str(img_path), type="png")

                # Add to PPTX
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(img_path),
                    left=Emu(0),
                    top=Emu(0),
                    width=slide_width,
                    height=slide_height,
                )

        browser.close()

    prs.save(str(output_path))
    print(f"\nSaved {total} slides to {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Screenshot-based PPTX exporter for the presentation"
    )
    parser.add_argument("--html", type=Path, default=DEFAULT_HTML,
                        help="HTML file to export")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT,
                        help="Output PPTX path")
    parser.add_argument("--width", type=int, default=1920,
                        help="Viewport width (default: 1920)")
    parser.add_argument("--height", type=int, default=1080,
                        help="Viewport height (default: 1080)")
    parser.add_argument("--live-only", action="store_true",
                        help="Only export live deck slides (skip reference)")
    args = parser.parse_args()

    export(args.html, args.output, args.width, args.height, args.live_only)


if __name__ == "__main__":
    main()
