#!/usr/bin/env python3
"""
Generate a branded PPTX deck from the Jina AI + EIS presentation content.

Usage:
    cd docs/presentation
    pip install -r requirements.txt
    python export_pptx.py
    # Output: output/jina-eis-101.pptx
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = Path(__file__).parent
OUTPUT_DIR = SCRIPT_DIR / "output"
ASSETS_DIR = SCRIPT_DIR / "assets"

ELASTIC_BLUE = RGBColor(0x0B, 0x64, 0xDD)
ELASTIC_TEAL = RGBColor(0x02, 0xBC, 0xB7)
ELASTIC_GOLD = RGBColor(0xFE, 0xC5, 0x14)
ELASTIC_PINK = RGBColor(0xF0, 0x4E, 0x98)
ELASTIC_GREEN = RGBColor(0x9A, 0xDC, 0x30)
ELASTIC_LIGHT_BLUE = RGBColor(0x1B, 0xA9, 0xF5)
JINA_TEAL = RGBColor(0x00, 0x91, 0x91)
JINA_RED = RGBColor(0xEB, 0x61, 0x61)
JINA_YELLOW = RGBColor(0xFB, 0xCB, 0x67)

BG_DARK = RGBColor(0x0F, 0x11, 0x17)
BG_CARD = RGBColor(0x1A, 0x1D, 0x27)
TEXT_WHITE = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DIM = RGBColor(0x64, 0x74, 0x8B)
SCENARIO_GOLD = RGBColor(0xD4, 0xA8, 0x3C)
CODE_BG = RGBColor(0x16, 0x19, 0x22)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, font_name="Calibri",
                 alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=TEXT_MUTED, bold_color=TEXT_WHITE, font_name="Calibri",
                  spacing=Pt(6)):
    """Add a text box with multiple lines. Lines starting with '**' are bold."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing

        is_bold = line.startswith("**")
        clean = line.strip("*").strip()

        if " -- " in clean and is_bold:
            parts = clean.split(" -- ", 1)
            run1 = p.add_run()
            run1.text = parts[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = bold_color
            run1.font.bold = True
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = " — " + parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = font_name
        else:
            run = p.add_run()
            run.text = clean
            run.font.size = Pt(font_size)
            run.font.color.rgb = bold_color if is_bold else color
            run.font.bold = is_bold
            run.font.name = font_name

    return txBox


def add_accent_bar(slide, left, top, width, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_section_label(slide, text, color, top=Inches(0.6)):
    add_text_box(slide, Inches(0.8), top, Inches(4), Inches(0.3),
                 text.upper(), font_size=10, color=color, bold=True)


def add_scenario_intro(slide, text, top=Inches(0.9)):
    """Add a scenario narrative intro with gold left-accent styling."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.8), top, Pt(3), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ELASTIC_GOLD
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1.0), top, Inches(11), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = SCENARIO_GOLD
    p.font.italic = True
    p.font.name = "Calibri"


def add_table(slide, left, top, width, rows_data, col_widths=None, font_size=11):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width,
                                          Inches(0.35 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = TEXT_MUTED
                else:
                    paragraph.font.color.rgb = TEXT_WHITE

            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = RGBColor(0x15, 0x17, 0x20)
            else:
                cell.fill.fore_color.rgb = BG_CARD if r % 2 == 1 else BG_DARK

    return table_shape


def add_code_box(slide, left, top, width, height, text, font_size=11):
    """Add a code block with monospace font on a dark background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x2A, 0x2E, 0x3A)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Courier New"
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(1)
    return shape


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    jina_logo = ASSETS_DIR / "jina-logo.png"
    elastic_logo = ASSETS_DIR / "elastic-logo.png"

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    if jina_logo.exists():
        slide.shapes.add_picture(str(jina_logo), Inches(4.5), Inches(1.2),
                                  height=Inches(0.8))
    add_text_box(slide, Inches(6.4), Inches(1.3), Inches(0.5), Inches(0.5),
                 "+", font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    if elastic_logo.exists():
        slide.shapes.add_picture(str(elastic_logo), Inches(7.2), Inches(1.1),
                                  height=Inches(1.0))

    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1),
                 '"We Might Be Illegal in Europe"',
                 font_size=40, color=TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.5), Inches(4.0), Inches(8), Inches(0.6),
                 "Building a compliance search tool with Jina AI + Elastic Inference Service "
                 "— before the lawyers' deadline.",
                 font_size=18, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.5), Inches(5.2), Inches(8), Inches(0.4),
                 "A story about shipping first and reading the fine print second  |  "
                 "~30 min talk + hands-on  |  Internal Enablement",
                 font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    add_accent_bar(slide, Inches(4), Inches(6.2), Inches(5.3), JINA_TEAL)

    # =========================================================================
    # SLIDE 2: What Is Jina AI
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Part 1 — Jina AI", JINA_TEAL)
    add_scenario_intro(slide, "Your company just shipped facial recognition for user onboarding. "
                       "Legal has questions. You need an emergency response team. Meet Jina AI.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8), Inches(0.7),
                 "What Is Jina AI?", font_size=32, color=TEXT_WHITE, bold=True)

    lines = [
        "**Berlin-based AI company, acquired by Elastic in October 2025",
        "Builds search foundation models — the ML models that sit between your data and your search results. They turn \"dumb\" keyword lookup into intelligent meaning-based retrieval.",
        "Not a search engine — they build the models that make search engines smart.",
        "**All current models (v3+): CC BY-NC 4.0 — free to evaluate, commercial use through Elastic",
    ]
    add_multiline(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.5),
                  lines, font_size=14)

    families = [
        "**Embeddings -- Text & images → vectors for semantic search",
        "**Rerankers -- Re-score and re-order results for precision",
        "**Reader -- Full service — URL/PDF → clean markdown or JSON",
        "**CLIP -- Cross-modal: find images with text, text with images",
        "**VLM -- Vision language model — look at an image, answer questions",
    ]
    add_text_box(slide, Inches(7.0), Inches(1.5), Inches(3), Inches(0.4),
                 "Five Product Families", font_size=16, color=JINA_TEAL, bold=True)
    add_multiline(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.5),
                  families, font_size=13, spacing=Pt(8))

    # =========================================================================
    # SLIDE 3: Key Concepts
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI", JINA_TEAL)
    add_scenario_intro(slide, "Before you can build anything, the PM wants to know what an "
                       "'embedding' is. Fair enough — you had to Google it too.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Key Concepts — Quick Primer", font_size=32,
                 color=TEXT_WHITE, bold=True)

    concepts = [
        ("What Is an Embedding?",
         "Text → list of numbers (a vector). Similar meanings = similar numbers. "
         "\"Fix a leaky faucet\" matches \"repairing dripping taps\" — that's semantic search."),
        ("What Is a Context Window?",
         "Max text a model processes in one pass, in tokens (~1 token ≈ 1 word). "
         "512 tokens ≈ 1 paragraph. 8K ≈ 6 pages. 32K ≈ 25 pages. Longer = fewer chunks, better quality."),
        ("What Is Matryoshka Representation Learning?",
         "Named after Russian nesting dolls. Front-loads the most important info into "
         "the first dimensions. Truncate 1,024 → 256 dims, keep 90%+ quality. Trade tiny accuracy loss for big storage savings."),
        ("What Are LoRA Adapters?",
         "Small plugin modules attached to a base model. One model, switch tasks "
         "(retrieval, classification, clustering) at query time. Deploy once, not three times."),
    ]

    for i, (term, definition) in enumerate(concepts):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.0)
        y = Inches(2.0) + row * Inches(2.3)

        add_text_box(slide, x, y, Inches(5.5), Inches(0.35),
                     term, font_size=13, color=JINA_TEAL, bold=True)
        add_text_box(slide, x, y + Inches(0.4), Inches(5.5), Inches(1.6),
                     definition, font_size=11, color=TEXT_MUTED)

    # =========================================================================
    # SLIDE 3b: Embedding Similarity Demo (Static)
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI — Key Concepts", JINA_TEAL)
    add_scenario_intro(slide, 'This is why "facial recognition" finds "biometric identification" '
                       '— same neighborhood in vector space.')
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Embedding Similarity — How It Works", font_size=28,
                 color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(0.35),
                 "Similar meanings cluster together in vector space — even with completely different words.",
                 font_size=12, color=TEXT_DIM)

    sim_headers = ["Cosine Sim", "Faucet", "Taps", "Plumbing", "Quantum", "Neural"]
    sim_rows = [
        sim_headers,
        ["Faucet", "—", "0.92", "0.89", "0.12", "0.18"],
        ["Taps", "0.92", "—", "0.87", "0.10", "0.15"],
        ["Plumbing", "0.89", "0.87", "—", "0.14", "0.21"],
        ["Quantum", "0.12", "0.10", "0.14", "—", "0.31"],
        ["Neural", "0.18", "0.15", "0.21", "0.31", "—"],
    ]
    add_table(slide, Inches(0.8), Inches(2.3), Inches(6.5), sim_rows, font_size=11)

    add_text_box(slide, Inches(8.0), Inches(2.3), Inches(4.5), Inches(0.3),
                 "2D Vector Space (conceptual)", font_size=12,
                 color=TEXT_DIM, bold=True, alignment=PP_ALIGN.CENTER)

    cluster_items = [
        ("Plumbing cluster (teal)", "\"fix a leaky faucet\", \"repairing dripping taps\", "
         "\"plumbing repair guide\" — all score 0.87–0.92 against each other"),
        ("Physics outlier (gold)", "\"quantum physics research\" — scores 0.10–0.14 against "
         "plumbing, very far away in vector space"),
        ("ML outlier (blue)", "\"train a neural network\" — scores 0.15–0.21 against plumbing, "
         "0.31 against quantum (slight topical proximity via science)"),
    ]
    for i, (label, desc) in enumerate(cluster_items):
        y = Inches(2.8) + i * Inches(1.3)
        colors_list = [JINA_TEAL, ELASTIC_GOLD, ELASTIC_BLUE]
        add_text_box(slide, Inches(8.0), y, Inches(4.5), Inches(0.3),
                     label, font_size=11, color=colors_list[i], bold=True)
        add_text_box(slide, Inches(8.0), y + Inches(0.35), Inches(4.5), Inches(0.8),
                     desc, font_size=10, color=TEXT_MUTED)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "This is how semantic search works: \"fix a leaky faucet\" matches "
                 "\"repairing dripping taps\" because they live in the same neighborhood "
                 "of vector space — despite sharing zero keywords.",
                 font_size=12, color=JINA_TEAL)

    # =========================================================================
    # SLIDE 4: Embeddings — v5 Hero
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI — Embeddings", JINA_TEAL)
    add_scenario_intro(slide, 'The regulation calls it "real-time remote biometric identification systems." '
                       'Your old search calls that "no results found." Time for a better model.')
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Embedding Models — Lead with v5", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(2.3), Inches(11.5), ELASTIC_GOLD)

    add_text_box(slide, Inches(0.8), Inches(2.15), Inches(7), Inches(0.4),
                 "jina-embeddings-v5-text-small     NEW DEFAULT  |  ON EIS",
                 font_size=16, color=ELASTIC_GOLD, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(7), Inches(1.2),
                 "The recommended default embedding model for Elastic, replacing ELSER "
                 "as the go-to for new deployments. Best search quality of any multilingual "
                 "model under 1B parameters.\n\nIf you're recommending ELSER today, "
                 "recommend v5-text-small instead.",
                 font_size=13, color=TEXT_MUTED)

    stats = [("677M", "Params"), ("1,024", "Dims"), ("32K", "Context (25 pages)"),
             ("119+", "Languages"), ("4", "LoRA adapters")]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.8) + i * Inches(1.8)
        add_text_box(slide, x, Inches(4.2), Inches(1.6), Inches(0.3),
                     val, font_size=16, color=TEXT_WHITE, bold=True,
                     font_name="Courier New")
        add_text_box(slide, x, Inches(4.55), Inches(1.6), Inches(0.2),
                     label.upper(), font_size=8, color=TEXT_DIM, bold=True)

    others = [
        ("v5-text-nano", "239M params, 768 dims, 32K context. Perfect for edge, latency, budget."),
        ("v4 (Multimodal)", "3.8B params, 2,048 dims. Text + images through one model."),
    ]
    for i, (name, desc) in enumerate(others):
        x = Inches(0.8) + i * Inches(5.8)
        add_text_box(slide, x, Inches(5.3), Inches(5.3), Inches(0.3),
                     name, font_size=14, color=TEXT_WHITE, bold=True)
        add_text_box(slide, x, Inches(5.7), Inches(5.3), Inches(0.5),
                     desc, font_size=11, color=TEXT_MUTED)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "All Jina embeddings: Matryoshka dimensions, LoRA adapters, binary quantization.",
                 font_size=12, color=JINA_TEAL)

    # =========================================================================
    # SLIDE 5: v5 vs ELSER
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI", JINA_TEAL)
    add_scenario_intro(slide, "Your cluster still has ELSER configured from 2024. "
                       "Here's why that won't save you.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "v5-text-small vs ELSER — Why the Switch?", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
                 "ELSER brought learned sparse retrieval to Elastic and worked well for English keyword-expansion. "
                 "v5 is a generational leap: true multilingual semantic search with 64× longer context.",
                 font_size=14, color=TEXT_MUTED)

    headers = ["Dimension", "ELSER v2", "Jina v5-text-small"]
    rows = [
        headers,
        ["Type", "Sparse (learned sparse)", "Dense (vector)"],
        ["Languages", "English-focused", "119+ languages"],
        ["Context Window", "512 tokens (~1 paragraph)", "32K tokens (~25 pages)"],
        ["Dimensions", "Sparse (variable)", "1,024 (truncatable via Matryoshka)"],
        ["Long Documents", "Needs aggressive chunking", "25 pages in a single pass"],
        ["Task Adapters", "None", "4 LoRA adapters"],
        ["On EIS", "Yes", "Yes"],
        ["Status", "Legacy default", "NEW DEFAULT"],
    ]
    add_table(slide, Inches(0.8), Inches(2.5), Inches(11.5), rows, font_size=12)

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
                 "Bottom line: ELSER = English keyword expansion. v5-text-small = true multilingual "
                 "semantic understanding with 64× more context. Always recommend v5-text-small for new deployments.",
                 font_size=13, color=ELASTIC_GOLD, bold=True)

    # =========================================================================
    # SLIDE 6: Rerankers
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI — Rerankers", JINA_TEAL)
    add_scenario_intro(slide, "Semantic search returns five articles. Only one of them actually "
                       "answers your legal team's question. The reranker figures out which one.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Rerankers — The Precision Layer", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6),
                 "Why two stages? Embedding search is fast but compares pre-computed vectors. A reranker reads each "
                 "candidate alongside the query in full detail. Speed from embeddings (100 candidates in ms), "
                 "precision from reranking (best 10).",
                 font_size=13, color=TEXT_MUTED)

    add_accent_bar(slide, Inches(0.8), Inches(2.6), Inches(5), JINA_TEAL)
    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(5), Inches(0.4),
                 "Reranker v3  (Recommended)", font_size=18,
                 color=ELASTIC_GOLD, bold=True)
    v3_lines = [
        "**Listwise scoring -- query + up to 64 documents simultaneously",
        "Compares documents against each other, not just the query",
        "**600M params, 131K token context",
        "100+ languages, outperforms models with 2.5× more params",
    ]
    add_multiline(slide, Inches(0.8), Inches(3.2), Inches(5), Inches(2.5),
                  v3_lines, font_size=12)

    add_accent_bar(slide, Inches(7.0), Inches(2.6), Inches(5), TEXT_DIM)
    add_text_box(slide, Inches(7.0), Inches(2.7), Inches(5), Inches(0.4),
                 "Reranker v2", font_size=18, color=TEXT_MUTED, bold=True)
    v2_lines = [
        "**Pointwise scoring -- each document scored independently",
        "**278M params, 1K token context",
        "100+ languages, 6× faster than v1",
        "Bonus: function-calling + code search support",
    ]
    add_multiline(slide, Inches(7.0), Inches(3.2), Inches(5), Inches(2.5),
                  v2_lines, font_size=12)

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
                 "The analogy: Embedding search = scanning a library catalog to find the right shelf. "
                 "The reranker = a librarian who reads the actual pages and tells you "
                 "which books best answer your question.",
                 font_size=13, color=JINA_TEAL)

    # =========================================================================
    # SLIDE 6b: Reranker Demo (Static)
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI — Rerankers", JINA_TEAL)
    add_scenario_intro(slide, "Your legal team's first question. Let's see if search can answer it.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Demo: Before & After Reranking", font_size=28,
                 color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(0.35),
                 'Query: "Can law enforcement use facial recognition?" — EU AI Act corpus',
                 font_size=12, color=TEXT_DIM)

    naive_items = [
        ("#1  Art. 5 — Prohibited AI Practices", "0.81"),
        ("#2  Art. 52 — Transparency for AI Systems", "0.79"),
        ("#3  Art. 6 — Classification of High-Risk AI", "0.77"),
        ("#4  Art. 26 — Obligations of Deployers", "0.74"),
        ("#5  Art. 9 — Risk Management System", "0.71"),
    ]
    reranked_items = [
        ("#1  Art. 6 — Classification of High-Risk AI", "0.96", "↑2"),
        ("#2  Art. 5 — Prohibited AI Practices", "0.93", "↓1"),
        ("#3  Art. 26 — Obligations of Deployers", "0.89", "↑1"),
        ("#4  Art. 9 — Risk Management System", "0.85", "↑1"),
        ("#5  Art. 52 — Transparency for AI Systems", "0.62", "↓3"),
    ]

    add_accent_bar(slide, Inches(0.8), Inches(2.2), Inches(5.5), TEXT_DIM)
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.3),
                 "SEMANTIC SEARCH ONLY", font_size=10, color=TEXT_DIM, bold=True)
    for i, (title, score) in enumerate(naive_items):
        y = Inches(2.7) + i * Inches(0.65)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.8), y, Inches(5.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = RGBColor(0x2A, 0x2E, 0x3A)
        add_text_box(slide, Inches(0.95), y + Inches(0.08), Inches(4.5), Inches(0.35),
                     title, font_size=11, color=TEXT_MUTED)
        add_text_box(slide, Inches(5.5), y + Inches(0.08), Inches(0.7), Inches(0.35),
                     score, font_size=11, color=TEXT_DIM, font_name="Courier New",
                     alignment=PP_ALIGN.RIGHT)

    add_accent_bar(slide, Inches(7.0), Inches(2.2), Inches(5.5), JINA_TEAL)
    add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(0.3),
                 "AFTER RERANKING (JINA V3)", font_size=10, color=JINA_TEAL, bold=True)
    for i, (title, score, move) in enumerate(reranked_items):
        y = Inches(2.7) + i * Inches(0.65)
        is_top = (i == 0)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(7.0), y, Inches(5.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x0D, 0x2A, 0x2A) if is_top else BG_CARD
        shape.line.color.rgb = JINA_TEAL if is_top else RGBColor(0x2A, 0x2E, 0x3A)
        add_text_box(slide, Inches(7.15), y + Inches(0.08), Inches(4.0), Inches(0.35),
                     title, font_size=11,
                     color=TEXT_WHITE if is_top else TEXT_MUTED)
        add_text_box(slide, Inches(11.3), y + Inches(0.08), Inches(0.5), Inches(0.35),
                     score, font_size=11,
                     color=JINA_TEAL if float(score) > 0.85 else TEXT_DIM,
                     font_name="Courier New", alignment=PP_ALIGN.RIGHT)
        move_color = RGBColor(0x4C, 0xAF, 0x87) if move.startswith("↑") else RGBColor(0xE5, 0x73, 0x73)
        add_text_box(slide, Inches(11.9), y + Inches(0.08), Inches(0.5), Inches(0.35),
                     move, font_size=11, color=move_color,
                     alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
                 "Art. 6 jumps from #3 to #1 — it directly addresses biometric classification "
                 "including facial recognition. Art. 52 (general transparency) drops to #5.",
                 font_size=12, color=JINA_TEAL)

    # =========================================================================
    # SLIDE 7: Jina Reader Service
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI", JINA_TEAL)
    add_scenario_intro(slide, "Step one: be able to read the regulation. The EU published it as "
                       "an HTML page with 200 lines of navigation, cookie banners, and a sidebar. "
                       "Reader fixes that.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Jina Reader — Content Pipeline Service", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
                 "A full service (not just a model) for preprocessing raw web content into clean, "
                 "structured data. Every RAG pipeline starts with clean data — Reader ensures AI models "
                 "get high-quality input instead of noisy HTML.",
                 font_size=13, color=TEXT_MUTED)

    modes = [
        ("Reader API", "r.jina.ai/{url}",
         "Prepend to any URL. Returns clean Markdown. Handles HTML, PDFs, dynamic JS pages."),
        ("Search Mode", "s.jina.ai/{query}",
         "Web search → top 5 results already converted to clean Markdown."),
        ("ReaderLM v2 (Model)", "1.5B params, 512K context",
         "Deploy locally. HTML → Markdown/JSON with schema extraction. 29 languages."),
    ]
    for i, (title, subtitle, desc) in enumerate(modes):
        x = Inches(0.8) + i * Inches(4.0)
        color = [ELASTIC_TEAL, ELASTIC_BLUE, ELASTIC_GOLD][i]
        add_accent_bar(slide, x, Inches(2.7), Inches(3.5), color)
        add_text_box(slide, x, Inches(2.8), Inches(3.5), Inches(0.3),
                     title, font_size=15, color=TEXT_WHITE, bold=True)
        add_text_box(slide, x, Inches(3.15), Inches(3.5), Inches(0.25),
                     subtitle, font_size=10, color=color,
                     font_name="Courier New")
        add_text_box(slide, x, Inches(3.5), Inches(3.5), Inches(1.5),
                     desc, font_size=11, color=TEXT_MUTED)

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                 "Outperforms models 20× its size on HTML-to-Markdown quality benchmarks.",
                 font_size=12, color=JINA_TEAL)

    # =========================================================================
    # SLIDE 7b: Reader Before/After (Static)
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI — Reader", JINA_TEAL)
    add_scenario_intro(slide, "This is what the EU AI Act looks like when you download it. "
                       "Good luck with Ctrl+F.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "See It: Before & After Jina Reader", font_size=28,
                 color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(0.35),
                 "What a web page looks like before and after processing with r.jina.ai/",
                 font_size=12, color=TEXT_DIM)

    add_accent_bar(slide, Inches(0.8), Inches(2.2), Inches(5.5), TEXT_DIM)
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.3),
                 "RAW HTML SOURCE", font_size=10, color=TEXT_DIM, bold=True)
    raw_html = (
        '<!DOCTYPE html>\n'
        '<html lang="en">\n'
        '<head>\n'
        '  <title>EU AI Act</title>\n'
        '  <link rel="stylesheet" href="/css/main.css">\n'
        '  <script src="/js/analytics.js"></script>\n'
        '  <script src="/js/cookie-consent.js"></script>\n'
        '</head>\n'
        '<body>\n'
        '  <nav class="top-bar">...</nav>\n'
        '  <div class="cookie-banner">We use...</div>\n'
        '  <div class="sidebar"> <!-- 200 lines --> </div>\n'
        '  <main>\n'
        '    <h2>Article 5 - Prohibited AI</h2>\n'
        '    <p>The following AI practices\n'
        '    shall be prohibited:</p>\n'
        '  </main>\n'
        '  <footer> <!-- 80 lines --> </footer>\n'
        '  <script src="/js/tracking.js"></script>\n'
        '</body></html>'
    )
    add_code_box(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3.8),
                 raw_html, font_size=9)

    add_accent_bar(slide, Inches(7.0), Inches(2.2), Inches(5.5), JINA_TEAL)
    add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(0.3),
                 "AFTER JINA READER → CLEAN MARKDOWN", font_size=10,
                 color=JINA_TEAL, bold=True)
    clean_md = (
        '# Article 5 – Prohibited AI Practices\n'
        '\n'
        'The following AI practices shall\n'
        'be prohibited:\n'
        '\n'
        '## (a) Subliminal Techniques\n'
        'The placing on the market, the putting\n'
        'into service, or the use of an AI\n'
        'system that deploys subliminal\n'
        'techniques beyond a person\'s\n'
        'consciousness...\n'
        '\n'
        '## (b) Exploitation of Vulnerabilities\n'
        'AI systems that exploit vulnerabilities\n'
        'of a specific group of persons due to\n'
        'their age, disability, or social or\n'
        'economic situation.\n'
        '\n'
        '---\n'
        'Source: EU AI Act\n'
        'Retrieved via r.jina.ai'
    )
    add_code_box(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(3.8),
                 clean_md, font_size=9)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
                 "Scripts, navigation, ads, cookie banners, sidebars, footers — all removed. "
                 "Clean, structured content ready for embedding.",
                 font_size=12, color=JINA_TEAL)

    # =========================================================================
    # SLIDE 8: CLIP + VLM
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI", JINA_TEAL)
    add_scenario_intro(slide, 'Your CTO asks: "Is our AWS architecture compliant?" You upload '
                       "the diagram. The VLM reads it and flags Amazon Rekognition as a biometric system.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Visual Intelligence — CLIP & VLM", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(5.5), JINA_TEAL)
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.4),
                 "CLIP v2 — Cross-Modal Search", font_size=18,
                 color=TEXT_WHITE, bold=True)
    clip_lines = [
        "Text and image encoders producing embeddings in the same vector space",
        "\"Red sports car\" text query → finds matching images",
        "**865M params, 89 languages, 512×512 resolution (4× v1)",
        "Matryoshka: truncate 1,024 → 64 dims with <1% loss",
        "Lighter than v4 — choose for speed/cost",
    ]
    add_multiline(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3),
                  clip_lines, font_size=12)

    add_accent_bar(slide, Inches(7.0), Inches(2.0), Inches(5.5), ELASTIC_BLUE)
    add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(0.4),
                 "VLM — Vision Language Model", font_size=18,
                 color=TEXT_WHITE, bold=True)
    vlm_lines = [
        "Generative: looks at an image, produces text answers",
        "\"What does this chart show?\" \"What text is in this scan?\"",
        "**2.4B params, 29 languages",
        "Best multilingual VQA among open 2B models",
        "Strong OCR (778/1000), lowest hallucination (90.3 POPE)",
    ]
    add_multiline(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(3),
                  vlm_lines, font_size=12)

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
                 "CLIP v2 embeds images for search (lightweight). v4 embeds images + text for "
                 "multimodal search (heavier). VLM generates text about images (answers questions). "
                 "They can be used together.",
                 font_size=12, color=JINA_TEAL)

    # =========================================================================
    # SLIDE 8b: VLM Demo (Static)
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI — VLM", JINA_TEAL)
    add_scenario_intro(slide, 'The CTO\'s Slack message: "Are we using any of these banned AI services?"')
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "VLM in Action: Image → Text Analysis", font_size=28,
                 color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(0.35),
                 "Upload an architecture diagram and the VLM tells you what it sees.",
                 font_size=12, color=TEXT_DIM)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.8), Inches(2.3), Inches(4.0), Inches(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape.line.color.rgb = RGBColor(0x2A, 0x2E, 0x3A)
    add_text_box(slide, Inches(1.2), Inches(2.8), Inches(3.2), Inches(1.8),
                 "[ Cloud Architecture Diagram ]\n\n"
                 "API Gateway → Lambda → Rekognition\n"
                 "             → S3 Bucket\n"
                 "Rekognition → DynamoDB → SNS Topic",
                 font_size=11, color=TEXT_DIM, font_name="Courier New")
    add_text_box(slide, Inches(1.2), Inches(4.6), Inches(3.2), Inches(0.35),
                 "cloud-architecture-v2.png", font_size=9, color=TEXT_DIM,
                 alignment=PP_ALIGN.CENTER)

    prompt_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(5.5), Inches(2.3), Inches(7.0), Inches(0.6))
    prompt_shape.fill.solid()
    prompt_shape.fill.fore_color.rgb = RGBColor(0x0D, 0x2A, 0x2A)
    prompt_shape.line.color.rgb = JINA_TEAL
    add_text_box(slide, Inches(5.6), Inches(2.35), Inches(1.0), Inches(0.2),
                 "PROMPT", font_size=8, color=JINA_TEAL, bold=True)
    add_text_box(slide, Inches(5.6), Inches(2.55), Inches(6.7), Inches(0.3),
                 "What does this architecture diagram show? Identify all services and their connections.",
                 font_size=10, color=TEXT_WHITE)

    add_text_box(slide, Inches(5.5), Inches(3.1), Inches(0.8), Inches(0.25),
                 "RESPONSE", font_size=8, color=JINA_TEAL, bold=True)
    vlm_response = (
        "This architecture shows a serverless image processing pipeline on AWS:\n\n"
        "• API Gateway — Entry point for client requests\n"
        "• Lambda — Orchestrates the image analysis workflow\n"
        "• Amazon Rekognition — AI/ML for image classification\n"
        "• S3 Bucket — Object storage for images and results\n"
        "• DynamoDB — NoSQL database for analysis metadata\n"
        "• SNS Topic — Notifications on completion\n\n"
        "⚠ If Rekognition is used for facial recognition,\n"
        "this system may be classified as high-risk AI under\n"
        "EU AI Act Article 6 / Annex III."
    )
    add_code_box(slide, Inches(5.5), Inches(3.4), Inches(7.0), Inches(3.3),
                 vlm_response, font_size=9)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                 "Image in → Text analysis out. The VLM identifies services, data flows, "
                 "and compliance implications automatically.",
                 font_size=12, color=JINA_TEAL)

    # =========================================================================
    # SLIDE 9: Quick Reference
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI", JINA_TEAL)
    add_scenario_intro(slide, "Your team's cheat sheet. Bookmark this one — "
                       "you'll need it during the hands-on.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Quick Reference — Every Current Model", font_size=32,
                 color=TEXT_WHITE, bold=True)

    headers = ["Model", "Type", "Params", "Context", "Langs", "EIS?", "Best For"]
    rows = [
        headers,
        ["v5-text-small", "Embedding", "677M", "32K", "119+", "Yes", "Default text search & RAG"],
        ["v5-text-nano", "Embedding", "239M", "32K", "Multi", "Yes", "Edge, low-latency, budget"],
        ["v4", "Embedding", "3.8B", "32K", "30+", "No", "Visual docs, multimodal"],
        ["jina-clip-v2", "CLIP", "865M", "8K", "89", "No", "Cross-modal text ↔ image"],
        ["reranker-v3", "Reranker", "600M", "131K", "100+", "Yes", "Best accuracy, listwise"],
        ["reranker-v2", "Reranker", "278M", "1K", "100+", "Yes", "Fast, function-calling"],
        ["ReaderLM v2", "SLM", "1.5B", "512K", "29", "No", "HTML/PDF → Markdown/JSON"],
        ["jina-vlm", "VLM", "2.4B", "Multi", "29", "No", "Visual QA, OCR"],
    ]
    add_table(slide, Inches(0.5), Inches(1.9), Inches(12.3), rows, font_size=11)

    # =========================================================================
    # SLIDE 10: Why It Matters
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI", JINA_TEAL)
    add_scenario_intro(slide, "Legal says the regulation text can't leave the cluster. "
                       "So no, you can't just call OpenAI.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Why This Matters", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4),
                 "For Customer Conversations", font_size=16,
                 color=JINA_TEAL, bold=True)
    scope_lines = [
        "**Pure text search? -- v5-text-small on EIS. Done.",
        "**Documents with images/charts? -- v4.",
        "**Tight on budget/latency? -- v5-text-nano.",
        "**Need reranking precision? -- Add Reranker v3.",
        "**Quick PoC? -- v5-text-small on EIS. Under an hour.",
    ]
    add_multiline(slide, Inches(0.8), Inches(2.6), Inches(5.2), Inches(3.5),
                  scope_lines, font_size=13)

    add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5), Inches(0.4),
                 "vs. OpenAI / Cohere", font_size=16,
                 color=JINA_TEAL, bold=True)
    comp_lines = [
        "**No external API calls -- runs inside Elasticsearch on EIS",
        "**No data leaving the cluster -- text transmitted to EIS only",
        "**Predictable costs -- part of Elastic Cloud subscription",
        "**Equal or better quality -- matches OpenAI on benchmarks",
        "**OpenAI-compatible API -- migration is straightforward",
    ]
    add_multiline(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(3.5),
                  comp_lines, font_size=13)

    # =========================================================================
    # SLIDE 11: Jina API Service
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Jina AI", JINA_TEAL)
    add_scenario_intro(slide, "But first — you prototype with the Jina API to prove the approach "
                       "works before touching production. No cluster needed.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Jina API Service — jina.ai", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
                 "All Jina models available directly via REST API at jina.ai — no Elastic cluster required. "
                 "Instant access, no credit card or registration. SOC 2 Type 1 & 2 compliant. "
                 "Great for prototyping, evaluation, or non-Elastic workflows.",
                 font_size=13, color=TEXT_MUTED)

    api_modes = [
        ("Reader", "r.jina.ai / s.jina.ai",
         "r.jina.ai/{url} → clean Markdown.\n"
         "s.jina.ai/{query} → web search → clean results.\n"
         "No API key needed for basic use."),
        ("Embeddings", "api.jina.ai/v1/embeddings",
         "Generate embeddings via REST.\n"
         "All Jina embedding models (v3, v5, CLIP).\n"
         "Pass text or images, get vectors back."),
        ("Reranker", "api.jina.ai/v1/rerank",
         "Pass query + documents.\n"
         "Get re-scored and re-ordered results.\n"
         "Same v2/v3 models available on EIS."),
    ]
    for i, (title, endpoint, desc) in enumerate(api_modes):
        x = Inches(0.8) + i * Inches(4.0)
        color = [ELASTIC_TEAL, ELASTIC_BLUE, ELASTIC_GOLD][i]
        add_accent_bar(slide, x, Inches(2.7), Inches(3.5), color)
        add_text_box(slide, x, Inches(2.8), Inches(3.5), Inches(0.3),
                     title, font_size=15, color=TEXT_WHITE, bold=True)
        add_text_box(slide, x, Inches(3.15), Inches(3.5), Inches(0.25),
                     endpoint, font_size=9, color=color,
                     font_name="Courier New")
        add_text_box(slide, x, Inches(3.5), Inches(3.5), Inches(1.5),
                     desc, font_size=11, color=TEXT_MUTED)

    add_text_box(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.8),
                 "MCP Server: Add mcp.jina.ai as an MCP server to access "
                 "Reader, Embeddings, and Reranker directly from LLMs and agents.",
                 font_size=12, color=JINA_TEAL)

    add_text_box(slide, Inches(7.0), Inches(5.3), Inches(5.5), Inches(0.8),
                 "API Key: jina.ai → \"How to get my API key?\" → API Key & Billing tab. "
                 "Free tier with rate limits; add key for higher throughput.",
                 font_size=12, color=ELASTIC_GOLD)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "Jina API = prototyping & non-Elastic use.  EIS = production Elastic deployments "
                 "(same models, managed inside Elasticsearch, no external calls).",
                 font_size=12, color=ELASTIC_LIGHT_BLUE, bold=True)

    # =========================================================================
    # SLIDE 12: What Is EIS
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, ELASTIC_BLUE)
    add_section_label(slide, "Part 2 — Elastic Inference Service", ELASTIC_BLUE)
    add_scenario_intro(slide, "Prototype works. Now make it production-ready without "
                       "provisioning GPU nodes at 2am.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "What Is EIS?", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.5),
                 "Historically, running ML models in Elasticsearch meant provisioning dedicated ML nodes, "
                 "managing GPU resources, scaling allocations, and monitoring model health. EIS eliminates all of that.",
                 font_size=13, color=TEXT_MUTED)

    eis_lines = [
        "**Managed ML inference as a service -- embeddings, reranking, and LLM chat without ML nodes",
        "**No GPU management -- Elastic runs models on their infrastructure",
        "**Billed per million tokens -- not VCU hours or node count",
        "**Data stays in your cluster -- only text sent for inference is transmitted, discarded after processing",
    ]
    add_multiline(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3),
                  eis_lines, font_size=14)

    add_text_box(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(2.5),
                 "The key idea: Instead of provisioning and scaling ML nodes, EIS handles "
                 "all of that. Configure an inference endpoint — one API call — and "
                 "Elasticsearch uses it automatically for ingest, search, and chat. "
                 "Zero infrastructure to manage.",
                 font_size=14, color=ELASTIC_LIGHT_BLUE)

    # =========================================================================
    # SLIDE 12: Where Available
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "EIS", ELASTIC_BLUE)
    add_scenario_intro(slide, "Good news: your customer's cluster already has access. "
                       "Here's the deployment map.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Where Is EIS Available?", font_size=32,
                 color=TEXT_WHITE, bold=True)

    headers = ["Deployment Type", "Availability", "Setup Required", "Version"]
    rows = [
        headers,
        ["Elastic Cloud Serverless", "Built-in, zero setup", "None — pre-configured", "GA"],
        ["Elastic Cloud Hosted", "Built-in, zero setup", "None — pre-configured", "GA"],
        ["Self-managed (ECE, ECK, standalone)", "Via Cloud Connect", "Connect to Cloud, enable EIS", "9.3+"],
    ]
    add_table(slide, Inches(0.8), Inches(2.0), Inches(11.5), rows, font_size=13)

    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.5),
                 "Cloud Connect for self-managed: Bridges your cluster to Elastic Cloud services "
                 "without hosting them locally. After enabling, Elasticsearch automatically creates "
                 "inference endpoints and Kibana AI connectors. Requires Enterprise license or active "
                 "trial. Stack 9.3+.",
                 font_size=13, color=ELASTIC_LIGHT_BLUE)

    # =========================================================================
    # SLIDE 13: Three Modes
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "EIS", ELASTIC_BLUE)
    add_scenario_intro(slide, "Three ways to wire it up. For the compliance tool, Mode 1 is "
                       "all you need — the models are already there.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "The Three Modes", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(3.5), ELASTIC_TEAL)
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(3.5), Inches(0.3),
                 "MODE 1: PRE-CONFIGURED (DEFAULT)", font_size=10,
                 color=ELASTIC_TEAL, bold=True)
    add_code_box(slide, Inches(0.8), Inches(2.5), Inches(3.5), Inches(2.2),
                 'Exist out of the box. No API calls.\n\n'
                 'Dot-prefixed IDs:\n'
                 '".jina-embeddings-v5-text-small"\n'
                 '".jina-embeddings-v5-text-nano"\n'
                 '".jina-reranker-v3"\n'
                 '".elser-2-elastic"',
                 font_size=10)

    add_accent_bar(slide, Inches(4.8), Inches(2.0), Inches(3.5), ELASTIC_BLUE)
    add_text_box(slide, Inches(4.8), Inches(2.1), Inches(3.5), Inches(0.3),
                 "MODE 2: CUSTOM EIS ENDPOINTS", font_size=10,
                 color=ELASTIC_LIGHT_BLUE, bold=True)
    add_code_box(slide, Inches(4.8), Inches(2.5), Inches(3.5), Inches(2.2),
                 'PUT _inference/text_embedding/my-v5\n'
                 '{\n'
                 '  "service": "elastic",\n'
                 '  "service_settings": {\n'
                 '    "model_id":\n'
                 '      "jina-embeddings-v5-text-small"\n'
                 '  }\n'
                 '}',
                 font_size=10)

    add_accent_bar(slide, Inches(8.8), Inches(2.0), Inches(3.5), ELASTIC_GOLD)
    add_text_box(slide, Inches(8.8), Inches(2.1), Inches(3.5), Inches(0.3),
                 "MODE 3: THIRD-PARTY / BYOK", font_size=10,
                 color=ELASTIC_GOLD, bold=True)
    add_code_box(slide, Inches(8.8), Inches(2.5), Inches(3.5), Inches(2.2),
                 'PUT _inference/text_embedding/my-oai\n'
                 '{\n'
                 '  "service": "openai",\n'
                 '  "service_settings": {\n'
                 '    "model_id":\n'
                 '      "text-embedding-3-small",\n'
                 '    "api_key": "<YOUR_KEY>"\n'
                 '  }\n'
                 '}',
                 font_size=10)

    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.8),
                 "Modes 1 & 2 run on Elastic's infrastructure (EIS). "
                 "Mode 3 sends requests to external providers. "
                 "All three use the same Inference API.",
                 font_size=14, color=ELASTIC_LIGHT_BLUE, bold=True)

    # =========================================================================
    # SLIDE 14: Models on EIS — LLMs
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "EIS", ELASTIC_BLUE)
    add_scenario_intro(slide, "The full menu. For our compliance tool we need embeddings, "
                       "a reranker, and an LLM for the agent. All here.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "What's Available on EIS? — LLM Chat Models", font_size=28,
                 color=TEXT_WHITE, bold=True)

    headers = ["Author", "Model", "ID", "Status"]
    rows = [
        headers,
        ["Anthropic", "Claude Opus 4.6", "anthropic-claude-4.6-opus", "GA"],
        ["Anthropic", "Claude Sonnet 4.6", "anthropic-claude-4.6-sonnet", "Beta"],
        ["Anthropic", "Claude Opus 4.5", "anthropic-claude-4.5-opus", "GA"],
        ["Anthropic", "Claude Sonnet 4.5", "anthropic-claude-4.5-sonnet", "GA"],
        ["Anthropic", "Claude Haiku 4.5", "anthropic-claude-4.5-haiku", "Beta"],
        ["Anthropic", "Claude Sonnet 3.7", "rainbow-sprinkles", "GA (legacy)"],
        ["Google", "Gemini 2.5 Pro", "google-gemini-2.5-pro", "GA"],
        ["Google", "Gemini 2.5 Flash", "google-gemini-2.5-flash", "GA"],
        ["Google", "Gemini 2.5 Flash Lite", "google-gemini-2.5-flash-lite", "Beta"],
        ["OpenAI", "GPT-5.2", "openai-gpt-5.2", "GA"],
        ["OpenAI", "GPT-4.1", "openai-gpt-4.1", "GA"],
        ["OpenAI", "GPT-4.1 Mini", "openai-gpt-4.1-mini", "GA"],
        ["OpenAI", "GPT-OSS 120B", "openai-gpt-oss-120b", "GA"],
    ]
    add_table(slide, Inches(0.8), Inches(1.9), Inches(11.5), rows, font_size=10)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.3),
                 "All LLMs: 0-day data retention  |  Data NOT used to train models  |  US region",
                 font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 15: Models on EIS — Embeddings + Rerankers
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "EIS", ELASTIC_BLUE)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.7),
                 "What's Available on EIS? — Embeddings & Rerankers", font_size=28,
                 color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(0.3),
                 "Embedding Models", font_size=14, color=JINA_TEAL, bold=True)
    emb_rows = [
        ["Model", "ID", "Dims", "Service"],
        ["Jina v5 Small", "jina-embeddings-v5-text-small", "1,024", "EIS"],
        ["Jina v5 Nano", "jina-embeddings-v5-text-nano", "768", "EIS"],
        ["Jina v3", "jina-embeddings-v3", "1,024", "EIS"],
        ["ELSER v2 (EIS)", "elser_model_2", "sparse", "EIS"],
        ["ELSER v2 (ML node)", ".elser_model_2_linux-x86_64", "sparse", "ES"],
        ["E5 Large", "microsoft-multilingual-e5-large", "1,024", "EIS (Beta)"],
        ["E5 Small (ML)", ".multilingual-e5-small_linux-x86_64", "—", "ES"],
        ["Gemini Emb 001", "google-gemini-embedding-001", "768", "EIS"],
        ["OpenAI Emb 3 Large", "openai-text-embedding-3-large", "3,072", "EIS"],
        ["OpenAI Emb 3 Small", "openai-text-embedding-3-small", "1,536", "EIS"],
    ]
    add_table(slide, Inches(0.8), Inches(2.3), Inches(7.0), emb_rows, font_size=9)

    add_text_box(slide, Inches(8.5), Inches(1.9), Inches(3), Inches(0.3),
                 "Rerankers", font_size=14, color=JINA_TEAL, bold=True)
    rerank_rows = [
        ["Model", "ID", "Service"],
        ["Jina Reranker v3", "jina-reranker-v3", "EIS"],
        ["Jina Reranker v2", "jina-reranker-v2-base-multilingual", "EIS"],
        ["Rerank v1 (ML)", ".rerank-v1", "ES"],
    ]
    add_table(slide, Inches(8.5), Inches(2.3), Inches(4.0), rerank_rows, font_size=9)

    # =========================================================================
    # SLIDE 16: How It Works
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "EIS", ELASTIC_BLUE)
    add_scenario_intro(slide, "Three API calls. That's it. Map a field, search it, rerank the results. "
                       "This is what you're about to build.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "How It Works — End to End", font_size=32,
                 color=TEXT_WHITE, bold=True)

    steps = [
        ("Step 1", "Map", "semantic_text field"),
        ("Step 2", "Index", "Auto-embed on ingest"),
        ("Step 3", "Search", "match → semantic"),
        ("Step 4", "Rerank", "text_similarity_reranker"),
        ("Step 5", "Results", "Precision-ranked hits"),
    ]
    step_w = Inches(2.0)
    arrow_w = Inches(0.4)
    start_x = Inches(0.8)
    for i, (num, title, detail) in enumerate(steps):
        x = start_x + i * (step_w + arrow_w)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, Inches(2.2), step_w, Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = RGBColor(0x2A, 0x2E, 0x3A)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(8)
        p.font.color.rgb = TEXT_DIM
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = detail
        p3.font.size = Pt(9)
        p3.font.color.rgb = TEXT_MUTED
        p3.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            add_text_box(slide, x + step_w, Inches(2.6),
                         arrow_w, Inches(0.4), "→",
                         font_size=18, color=TEXT_DIM,
                         alignment=PP_ALIGN.CENTER)

    add_code_box(slide, Inches(0.5), Inches(3.9), Inches(3.9), Inches(3.0),
                 'PUT my-index\n'
                 '{\n'
                 '  "mappings": {\n'
                 '    "properties": {\n'
                 '      "text": {\n'
                 '        "type": "semantic_text",\n'
                 '        "inference_id":\n'
                 '          ".jina-embeddings-v5-\n'
                 '           text-small"\n'
                 '      }\n'
                 '    }\n'
                 '  }\n'
                 '}',
                 font_size=9)

    add_code_box(slide, Inches(4.7), Inches(3.9), Inches(3.9), Inches(3.0),
                 'GET my-index/_search\n'
                 '{\n'
                 '  "query": {\n'
                 '    "match": {\n'
                 '      "text": "soccer"\n'
                 '    }\n'
                 '  }\n'
                 '}\n'
                 '\n'
                 '// Returns "Aberdeen Football\n'
                 '// Club" — semantic match',
                 font_size=9)

    add_code_box(slide, Inches(8.9), Inches(3.9), Inches(3.9), Inches(3.0),
                 'GET my-index/_search\n'
                 '{\n'
                 '  "retriever": {\n'
                 '    "text_similarity_reranker":{\n'
                 '      "retriever": {\n'
                 '        "standard": {\n'
                 '          "query": {\n'
                 '            "match": {\n'
                 '              "text": "soccer"\n'
                 '      }}}},\n'
                 '      "inference_id":\n'
                 '        ".jina-reranker-v3",\n'
                 '      "field": "text"\n'
                 '  }}\n'
                 '}',
                 font_size=9)

    # =========================================================================
    # SLIDE 17: Hands-On
    # =========================================================================
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_section_label(slide, "Part 3 — Hands-On", ELASTIC_GOLD)
    add_scenario_intro(slide, "Your turn. Build the tool that would have saved your team "
                       "a week of panic.")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.7),
                 "Let's Build It", font_size=32,
                 color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(10), Inches(0.5),
                 "Three notebooks, one full pipeline — from raw PDF to "
                 "precision-ranked semantic search. Then: the interactive demo app.",
                 font_size=16, color=TEXT_MUTED)

    notebooks = [
        ("01", "Ingest", "Jina Reader → EU AI Act PDF → clean markdown, chunked by article.",
         JINA_TEAL),
        ("02", "Index", "Jina Embeddings v5 on EIS. semantic_text — zero-config vector search.",
         ELASTIC_BLUE),
        ("03", "Rerank", "Jina Reranker via text_similarity_reranker. Before vs after comparison.",
         ELASTIC_GOLD),
    ]
    for i, (num, title, desc, color) in enumerate(notebooks):
        x = Inches(0.8) + i * Inches(4.0)
        add_accent_bar(slide, x, Inches(3.0), Inches(3.5), color)
        add_text_box(slide, x, Inches(3.1), Inches(0.6), Inches(0.6),
                     num, font_size=28, color=color, bold=True)
        add_text_box(slide, x + Inches(0.6), Inches(3.2), Inches(2.8), Inches(0.4),
                     title, font_size=18, color=TEXT_WHITE, bold=True)
        add_text_box(slide, x, Inches(3.8), Inches(3.5), Inches(1.5),
                     desc, font_size=12, color=TEXT_MUTED)

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1),
                 "Then: The Interactive Demo — Search tab (semantic search + reranking)  |  "
                 "Agent tab (AI agent + VLM diagram analysis)  |  "
                 "Data Lab (live ingestion pipeline)",
                 font_size=14, color=TEXT_MUTED, bold=True)

    # =========================================================================
    # Save
    # =========================================================================
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUTPUT_DIR / "jina-eis-101.pptx"
    prs.save(str(output_path))
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
